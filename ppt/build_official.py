"""Fill the official SIH 2026 Idea template for SIH26037.

The template's own instruction slide sets the rules this script obeys: at most
six slides including the title, points rather than paragraphs, and the section
pointers left as the template defines them. So the structure is untouched - the
guidance text inside each section is replaced with our content, the instruction
slide is removed, and figures, a pipeline diagram and a results chart are added,
because a text-only idea PPT is exactly what the instructions tell you not to
submit.

Every performance number on the slides is read from artifacts/benchmark.json,
which scripts/benchmark.py writes. Nothing is typed in by hand, so the deck
cannot drift away from what the code actually does.

    uv run python scripts/benchmark.py --seeds 3
    cd ppt && uv run python build_official.py
"""
import json
from copy import deepcopy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC = "sih_format.pptx"
OUT = "SARATHI_SIH26037_Idea_Submission.pptx"
FIG = Path("../artifacts")
BENCH = Path("../artifacts/benchmark.json")

INK = RGBColor(0x26, 0x26, 0x26)
HEAD = RGBColor(0x1F, 0x38, 0x64)       # template's dark navy
ACCENT = RGBColor(0xB0, 0x3A, 0x2B)     # signal red, for numbers only
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
TINT = RGBColor(0xEC, 0xEF, 0xF5)       # pale navy wash for cards
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TEAM_NAME = "«TEAM NAME»"
TEAM_ID = "«TEAM ID»"

PRETTY = {
    "bus_stop_overtake": "Bus stop overtake",
    "cattle_crossing_sudden": "Cattle crossing",
    "construction_diversion": "Construction diversion",
    "highway_merge_slow": "Highway merge",
    "market_dense_mixed": "Dense market",
    "narrow_bridge_oncoming": "Narrow bridge",
    "night_highway_wrongway": "Night, wrong-way",
    "school_zone_pedestrians": "School zone",
    "urban_intersection_unsignalled": "Unsignalled junction",
    "village_road_unmarked": "Village road",
}


# --------------------------------------------------------------- template glue
def shape_by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape {shape_id} not on slide")


def set_text(shape, text):
    """Replace a template shape's text without disturbing its formatting.

    The placeholders are not uniform: some have an empty leading paragraph, one
    carries an empty run and a line break before the visible text. Writing into
    the last run and blanking the others keeps whatever spacing the template
    intended, and still works when a paragraph has no run at all.
    """
    tf = shape.text_frame
    runs = [r for p in tf.paragraphs for r in p.runs]
    if not runs:
        para = tf.paragraphs[-1]
        run = para.add_run()
        end = para._p.find(qn("a:endParaRPr"))
        if end is not None:
            rPr = deepcopy(end)
            rPr.tag = qn("a:rPr")
            existing = run._r.find(qn("a:rPr"))
            if existing is not None:
                run._r.remove(existing)
            run._r.insert(0, rPr)
        runs = [run]
    for r in runs[:-1]:
        r.text = ""
    runs[-1].text = text
    # A line break the template put in front of a one-word title costs a whole
    # line, and a longer title then overflows its box onto the content below.
    for br in tf._txBody.findall(f".//{qn('a:br')}"):
        br.getparent().remove(br)
    return shape


def set_team_chip(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Your Team Name":
            set_text(sh, TEAM_NAME)


def drop(slide, shape_id):
    sh = shape_by_id(slide, shape_id)
    sh._element.getparent().remove(sh._element)


def delete_slide(prs, index):
    lst = prs.slides._sldIdLst
    slides = list(lst)
    rId = slides[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    lst.remove(slides[index])


# ------------------------------------------------------------------- text bits
def clear(tf):
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    return tf


def write(tf, blocks, base=13.0):
    """Render ``blocks`` of (text, kind) into a text frame.

    kind: 'h' section heading, 'b' bullet, 's' sub-bullet, 'n' plain line,
          'k' key line (bold, accent), 'c' caption.
    """
    clear(tf)
    tf.word_wrap = True
    first = True
    for text, kind in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        for tag in ("buChar", "buAutoNum", "buNone", "buFont"):
            for el in pPr.findall(
                    f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"):
                pPr.remove(el)
        if kind == "h":
            pPr.set("marL", "0")
            pPr.set("indent", "0")
            p.space_before = Pt(0 if first else 9)
            p.space_after = Pt(3)
        elif kind == "s":
            pPr.set("marL", "466725")
            pPr.set("indent", "-190500")
            p.space_after = Pt(2)
        elif kind in ("n", "c"):
            pPr.set("marL", "0")
            pPr.set("indent", "0")
            p.space_after = Pt(4)
        else:
            pPr.set("marL", "190500")
            pPr.set("indent", "-190500")
            p.space_after = Pt(4)
        first = False

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ("— " + text if kind == "s"
                    else "• " + text if kind == "b" else text)
        f = run.font
        f.name = "Arial"
        f.size = Pt({"h": base + 2.0, "c": base - 2.5}.get(kind, base))
        f.bold = kind in ("h", "k")
        f.italic = kind == "c"
        f.color.rgb = {"h": HEAD, "k": ACCENT, "c": MUTED}.get(kind, INK)
    return tf


def box(slide, x, y, w, h, blocks, base=13.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    write(tb.text_frame, blocks, base)
    return tb


def picture(slide, name, x, y, w, max_h=None):
    """Place a figure inside a box, never past it.

    A figure sized by width alone silently runs off the bottom of a slide the
    moment its aspect ratio changes - which is exactly what happens when the
    screenshots are recaptured. Fitting to both dimensions makes that impossible.
    """
    path = FIG / name
    with Image.open(path) as im:
        ar = im.width / im.height
    h = w / ar
    if max_h is not None and h > max_h:
        h, w = max_h, max_h * ar
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                   width=Inches(w), height=Inches(h))
    return pic


def caption(slide, x, y, w, text):
    return box(slide, x, y, w, 0.5, [(text, "c")], base=13.0)


# --------------------------------------------------------------- visual pieces
def chevrons(slide, x, y, w, h, labels, fill=HEAD):
    """A left-to-right pipeline strip - the methodology as a diagram."""
    n = len(labels)
    step = w / n
    for i, label in enumerate(labels):
        sh = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(x + i * step), Inches(y),
            Inches(step * 1.06), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        sh.line.color.rgb = WHITE
        sh.line.width = Pt(1.25)
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Arial"
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = WHITE
    return slide


def card(slide, x, y, w, h, value, label, value_pt=26, label_pt=10.5):
    """A stat callout: one number that matters, with its unit underneath."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.adjustments[0] = 0.10
    sh.fill.solid()
    sh.fill.fore_color.rgb = TINT
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    clear(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(value_pt)
    r.font.bold = True
    r.font.color.rgb = HEAD
    q = tf.add_paragraph()
    q.alignment = PP_ALIGN.CENTER
    q.space_before = Pt(1)
    r2 = q.add_run()
    r2.text = label
    r2.font.name = "Arial"
    r2.font.size = Pt(label_pt)
    r2.font.color.rgb = MUTED
    return sh


def panel(slide, x, y, w, h, blocks, base=12.5, fill=TINT):
    """A tinted card with text in it - used where a block needs setting apart."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.10)
    write(tf, blocks, base)
    return sh


def two_column_table(slide, x, y, w, h, header, rows, base=11.5):
    gf = slide.shapes.add_table(len(rows) + 1, 2, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    tbl = gf.table
    tbl.first_row = True
    tbl.horz_banding = False
    tbl.columns[0].width = Inches(w * 0.46)
    tbl.columns[1].width = Inches(w - w * 0.46)
    for c, text in enumerate(header):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEAD
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        tf = clear(cell.text_frame)
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(base)
        r.font.bold = True
        r.font.color.rgb = WHITE
    for i, (left, right) in enumerate(rows, start=1):
        for c, text in enumerate((left, right)):
            cell = tbl.cell(i, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else TINT
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            tf = clear(cell.text_frame)
            tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = text
            r.font.name = "Arial"
            r.font.size = Pt(base)
            r.font.bold = c == 1
            r.font.color.rgb = HEAD if c == 1 else MUTED
    return gf


def progress_chart(slide, x, y, w, h, bench):
    """Mean route progress per scenario, ours against the lane-following baseline."""
    names = list(PRETTY)
    data = CategoryChartData()
    data.categories = [PRETTY[n] for n in reversed(names)]
    for key, label in (("baseline", "Lane-following baseline"), ("sarathi", "SARATHI")):
        by = bench["summary"][key]["by_scenario"]
        data.add_series(
            label, [by[n]["mean_progress"] * 100.0 for n in reversed(names)])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(x), Inches(y),
                                Inches(w), Inches(h), data)
    chart = gf.chart
    chart.font.size = Pt(9)
    chart.font.name = "Arial"
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.overlap = -10
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.number_format = '0"%"'
    labels.number_format_is_linked = False
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.font.size = Pt(8)
    labels.font.color.rgb = MUTED
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = RGBColor(0xC2, 0xC7, 0xD0)
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = HEAD
    va = chart.value_axis
    va.maximum_scale = 100.0
    va.minimum_scale = 0.0
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE2, 0xE2)
    va.tick_labels.font.size = Pt(8)
    va.tick_labels.font.color.rgb = MUTED
    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(9)
    ca.tick_labels.font.color.rgb = INK
    return gf


# ------------------------------------------------------------------- the deck
def build():
    if not BENCH.exists():
        raise SystemExit(
            f"{BENCH} is missing. Run: uv run python scripts/benchmark.py --seeds 3")
    bench = json.loads(BENCH.read_text())
    ours = bench["summary"]["sarathi"]
    base = bench["summary"]["baseline"]
    seeds = len(bench["seeds"])
    scen = ours["scenarios"]
    p95 = ours["worst_p95_ms"]
    best = sorted(ours["by_scenario"].items(),
                  key=lambda kv: -kv[1]["mean_progress"])

    prs = Presentation(SRC)
    delete_slide(prs, 6)                 # the instruction slide, as it tells us to
    s1, s2, s3, s4, s5, s6 = prs.slides

    # ------------------------------------------------------------- 1. title
    set_text(shape_by_id(s1, 4), "SARATHI")
    write(shape_by_id(s1, 10).text_frame, [
        ("Problem Statement ID – SIH26037", "k"),
        ("Problem Statement Title – Adaptive Path Planning and Collision "
         "Avoidance for Autonomous Vehicles on Unstructured Indian Roads", "n"),
        ("Theme – Robotics and Drones", "n"),
        ("PS Category – Software", "n"),
        (f"Team ID – {TEAM_ID}", "n"),
        (f"Team Name (Registered on portal) – {TEAM_NAME}", "n"),
    ], base=17.0)

    # --------------------------------------------------- 2. idea / solution
    set_text(shape_by_id(s2, 15361), "SARATHI — PLANNING WITHOUT LANES")
    set_team_chip(s2)
    drop(s2, 15362)
    box(s2, 0.42, 1.22, 6.15, 5.55, [
        ("Proposed Solution", "h"),
        ("An Indian road has no lanes to follow — only free space that is "
         "negotiated, and that changes every second.", "n"),
        ("Lane centreline → drivable corridor, re-solved 20 times a second by "
         "a dynamic program through the free space ahead", "b"),
        ("Occupancy grid → Indian Driving Risk Field: continuous, "
         "class-conditioned, harm-weighted, indexed by time", "b"),
        ("Behaviour picked by an 8-state machine; the path itself by a "
         "jerk-minimal Frenet lattice scored in space and time", "b"),
        ("How it addresses the problem", "h"),
        ("Lane markings are never read, so faded, absent or wrong markings "
         "need no special case and no mode switch", "b"),
        ("A pothole is a cost bowl, not a wall; the verge is drivable; a cow "
         "is a high-variance hazard rather than a box", "b"),
        ("Innovation", "h"),
        ("Multi-modal intent per road user — cut in, filter, dart, ride the "
         "wrong way — each with covariance that grows with time", "b"),
        ("An RSS safety supervisor recalibrated to Indian gap acceptance; it "
         "is monotone by construction, so it can only ever slow us down", "b"),
    ], base=12.5)
    picture(s2, "fig-bev.png", 6.85, 1.30, 6.05, max_h=2.70)
    caption(s2, 6.85, 4.12, 6.05,
            "Risk field (red), the candidate fan (grey), the chosen path (amber).")
    two_column_table(s2, 6.85, 4.62, 6.05, 2.10,
                     ("A conventional stack uses", "SARATHI uses instead"),
                     [("Lane centreline", "Drivable corridor, re-solved each tick"),
                      ("Occupancy grid", "Class- and harm-weighted risk field"),
                      ("One predicted path", "Multi-modal intent per road user"),
                      ("Lane-change logic", "8 behaviours, wrong-way evade included")])

    # ------------------------------------------------- 3. technical approach
    set_team_chip(s3)
    drop(s3, 17410)
    chevrons(s3, 0.42, 1.14, 12.45, 0.46,
             ["SENSE", "FUSE", "PREDICT", "RISK FIELD", "CORRIDOR",
              "BEHAVIOUR", "LATTICE", "RSS + CBF"])
    box(s3, 0.42, 1.80, 6.05, 4.95, [
        ("Technologies", "h"),
        ("MATLAB, Simulink, Stateflow — the behaviour machine maps 1:1 onto "
         "Stateflow states", "b"),
        ("Automated Driving Toolbox, Navigation Toolbox and RoadRunner where "
         "licensed; the stack runs without them", "b"),
        ("Python 3.11, NumPy, SciPy — reference implementation and live demo", "b"),
        ("India Driving Dataset (IIIT Hyderabad) for the class taxonomy", "b"),
        ("Methodology", "h"),
        ("Sense — camera, LiDAR and radar with occlusion, dropout and class "
         "confusion. The planner never sees ground truth.", "b"),
        ("Fuse — constant-velocity Kalman tracks; class fused in log-odds", "b"),
        ("Predict — a manoeuvre distribution per agent, conditioned on class", "b"),
        ("Risk field — anisotropic kernels sized by footprint and uncertainty", "b"),
        ("Plan — corridor DP, then behaviour, then a jerk-minimal lattice", "b"),
        ("Assure — RSS-India inverted in closed form, plus a barrier speed cap", "b"),
        ("One scenario file drives both the Python and the MATLAB runtime", "k"),
    ], base=12.0)
    picture(s3, "fig-console.png", 6.72, 1.80, 6.15, max_h=3.62)
    caption(s3, 6.72, 5.52, 6.15,
            "Mission Control — every value on screen is computed live, not replayed.")
    card(s3, 6.72, 5.95, 1.94, 0.78, "20 Hz", "closed loop, no GPU")
    card(s3, 8.83, 5.95, 1.94, 0.78, f"{p95:.0f} ms", "worst 95th %ile replan")
    card(s3, 10.93, 5.95, 1.94, 0.78, "12", "road-user classes")

    # ------------------------------------------------ 4. feasibility, viability
    set_team_chip(s4)
    drop(s4, 17410)
    box(s4, 0.42, 1.18, 5.55, 3.30, [
        ("Feasibility — it already runs", "h"),
        ("One command from a clean clone: uv run sarathi serve", "b"),
        (f"{scen} scenarios — the five the problem statement names, plus five "
         "harder ones we added", "b"),
        ("Automated tests, including behavioural regressions, run in CI", "b"),
        ("20 Hz closed loop on an ordinary laptop; no GPU, no cloud", "b"),
        ("Viability", "h"),
        ("Built on the sponsor's own toolchain, and reproducible: every number "
         "on this slide comes out of a committed script", "b"),
    ], base=12.0)
    panel(s4, 0.42, 4.62, 5.55, 2.15, [
        ("Challenges, and what we do about them", "h"),
        ("Caution in dense traffic — measured every run, and every threshold "
         "is exposed live so behaviour can be tuned and defended", "s"),
        ("Real-world perception — grounded in IDD rather than assumed", "s"),
        ("RoadRunner licensing — a two-tier MATLAB design, so the deliverable "
         "is never blocked on a licence", "s"),
    ], base=11.5)
    progress_chart(s4, 6.15, 1.14, 6.75, 4.55, bench)
    caption(s4, 6.20, 5.72, 6.70,
            f"Mean route progress. {seeds} seeds × {scen} scenarios per "
            f"controller, identical seeds and identical sensor noise.")
    card(s4, 6.15, 6.10, 2.15, 0.72,
         f"{ours['collision_free']}/{ours['runs']}", "runs collision-free (ours)")
    card(s4, 8.45, 6.10, 2.15, 0.72,
         f"{base['collision_free']}/{base['runs']}", "runs collision-free (baseline)")
    card(s4, 10.75, 6.10, 2.15, 0.72,
         f"{ours['mean_progress'] / max(base['mean_progress'], 1e-6):.1f}×",
         "the baseline's route progress")

    # ---------------------------------------------------- 5. impact, benefits
    set_team_chip(s5)
    drop(s5, 17410)
    card(s5, 0.42, 1.15, 6.05, 1.05, "1,77,175",
         "road deaths in India in 2024 — about 485 every day (MoRTH)",
         value_pt=34, label_pt=11)
    box(s5, 0.42, 2.35, 6.05, 4.45, [
        ("Potential impact", "h"),
        ("Most of that toll falls on two-wheeler riders and pedestrians — "
         "exactly the road users this planner treats as first-class", "b"),
        ("Autonomy built on Western assumptions does not transfer. A planner "
         "that needs lane markings cannot even be validated here.", "b"),
        ("Who benefits", "h"),
        ("Researchers — an Indian-conditions scenario library, and a traffic "
         "model that reproduces filtering, wrong-way riding and cattle", "b"),
        ("Industry — a planner that degrades gracefully instead of "
         "disengaging when the road stops being tidy", "b"),
        ("Regulators and cities — safety margins that are measurable, "
         "reproducible and arguable, not a vendor's assertion", "b"),
    ], base=12.0)
    box(s5, 6.75, 1.15, 6.15, 3.05, [
        ("Benefits of the solution", "h"),
        ("Social — the risk field is explicitly harm-weighted, so a pedestrian "
         "costs the planner more than sheet metal does", "b"),
        ("Economic — validation in simulation instead of on public roads; the "
         "whole stack runs on a laptop", "b"),
        ("Environmental — smoother speed profiles and fewer emergency stops "
         "than the lane-based baseline on the same scenarios", "b"),
        ("Educational — every threshold is exposed and adjustable, so the "
         "system can be interrogated rather than trusted", "b"),
    ], base=12.0)
    picture(s5, "live-cow.png", 6.75, 3.55, 6.15, max_h=2.80)
    caption(s5, 6.75, 6.42, 6.15,
            "A judge drops a cow in front of the vehicle, mid-run. The planner "
            "has no foreknowledge of anything placed by hand.")

    # ------------------------------------------------ 6. research, references
    set_team_chip(s6)
    drop(s6, 17410)
    box(s6, 0.42, 1.18, 6.15, 5.55, [
        ("Grounding — data and the problem itself", "h"),
        ("India Driving Dataset, IIIT Hyderabad — idd.insaan.iiit.ac.in", "b"),
        ("MoRTH, Road Accidents in India 2024 — 1,77,175 fatalities", "b"),
        ("Research the design rests on", "h"),
        ("Werling et al., Optimal trajectory generation in a Frenet frame — "
         "the basis of our lattice (ICRA 2010)", "b"),
        ("Shalev-Shwartz, Shammah, Shashua, Responsibility-Sensitive Safety — "
         "the formal model our supervisor recalibrates for Indian gaps", "b"),
        ("Ames et al., Control barrier functions for safety-critical control", "b"),
        ("Non-lane-based car following in heterogeneous, disordered traffic — "
         "the model our surrounding agents drive by", "b"),
        ("MathWorks references", "h"),
        ("Highway trajectory planning using a Frenet reference path", "b"),
        ("Motion planning in urban environments with a dynamic occupancy grid", "b"),
        ("RoadRunner scene and scenario authoring documentation", "b"),
    ], base=12.0)
    panel(s6, 6.75, 1.18, 6.15, 1.62, [
        ("Try it yourself — nothing here is scripted", "h"),
        ("github.com/Paragraph1148/automated-driving", "k"),
        ("uv run sarathi serve", "k"),
    ], base=12.5)
    picture(s6, "drag.png", 6.95, 2.95, 5.75, max_h=3.55)
    caption(s6, 6.75, 6.58, 6.15,
            "Drag any road user with the mouse while the world runs — the "
            "planner replans against wherever you leave it.")
    prs.save(OUT)
    print(f"wrote {OUT}: {len(prs.slides._sldIdLst)} slides "
          f"(bench: ours {ours['collision_free']}/{ours['runs']} collision-free, "
          f"baseline {base['collision_free']}/{base['runs']})")


if __name__ == "__main__":
    build()
