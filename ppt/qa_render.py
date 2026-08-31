"""Render a .pptx to PNGs with PIL, for visual QA where LibreOffice cannot run.

soffice in this sandbox refuses every input ("source file could not be loaded"),
including a plain .txt, so the usual convert-to-pdf QA path is unavailable. This
draws the slides directly from the OOXML: shape fills, pictures, and wrapped text
at the resolved font size. Liberation Sans and Liberation Serif are metrically
compatible with Arial and Times New Roman, so text width - and therefore
overflow - is faithful even though the glyphs are not identical.

Usage: python qa_render.py deck.pptx [outdir] [--scale 120]
"""
import sys
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

FONTS = {
    ("sans", False, False): "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ("sans", True, False): "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ("sans", False, True): "/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf",
    ("sans", True, True): "/usr/share/fonts/truetype/liberation/LiberationSans-BoldItalic.ttf",
    ("serif", False, False): "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    ("serif", True, False): "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
    ("serif", False, True): "/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf",
    ("serif", True, True): "/usr/share/fonts/truetype/liberation/LiberationSerif-BoldItalic.ttf",
}
SERIFY = ("times", "serif", "cambria", "georgia", "bookman", "century", "garamond")
_cache = {}


def font(name, size_px, bold, italic):
    fam = "serif" if any(s in (name or "").lower() for s in SERIFY) else "sans"
    key = (fam, bool(bold), bool(italic), max(6, int(size_px)))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(FONTS[key[:3]], key[3])
    return _cache[key]


class Theme:
    """Resolve scheme colours (tx1, bg1, accent1, ...) to RGB via the master."""

    def __init__(self, prs):
        self.map = {}
        try:
            master = prs.slide_masters[0]
            theme = master.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
            scheme = theme._element.find(
                ".//" + qn("a:clrScheme"))
            for child in scheme:
                tag = child.tag.split("}")[1]
                srgb = child.find(qn("a:srgbClr"))
                sysc = child.find(qn("a:sysClr"))
                if srgb is not None:
                    self.map[tag] = srgb.get("val")
                elif sysc is not None:
                    self.map[tag] = sysc.get("lastClr", "000000")
            cmap = master._element.find(qn("p:clrMap"))
            if cmap is not None:
                self.alias = {k: v for k, v in cmap.attrib.items()}
            else:
                self.alias = {}
        except Exception:
            self.alias = {}

    def rgb(self, val):
        val = self.alias.get(val, val)
        val = {"dk1": "dk1", "lt1": "lt1", "dk2": "dk2", "lt2": "lt2"}.get(val, val)
        hexv = self.map.get(val)
        if hexv is None:
            hexv = {"tx1": "000000", "bg1": "FFFFFF", "tx2": "44546A",
                    "bg2": "E7E6E6"}.get(val, "808080")
        return tuple(int(hexv[i:i + 2], 16) for i in (0, 2, 4))


def solid_fill(fillel, theme):
    if fillel is None:
        return None
    solid = fillel.find(qn("a:solidFill"))
    if solid is None:
        return None
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is not None:
        v = srgb.get("val")
        return tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))
    sch = solid.find(qn("a:schemeClr"))
    if sch is not None:
        return theme.rgb(sch.get("val"))
    return None


def inherited_size(shape, level):
    """Walk slide -> layout -> master for a placeholder's default run size (pt)."""
    if not getattr(shape, "has_text_frame", False):
        return None
    def from_lststyle(txBody, lvl):
        if txBody is None:
            return None
        ls = txBody.find(qn("a:lstStyle"))
        if ls is None:
            return None
        el = ls.find(qn(f"a:lvl{lvl + 1}pPr"))
        if el is None:
            return None
        d = el.find(qn("a:defRPr"))
        return int(d.get("sz")) / 100 if d is not None and d.get("sz") else None

    sz = from_lststyle(shape.text_frame._txBody, level)
    if sz:
        return sz
    if not shape.is_placeholder:
        return None
    idx, ph_type = shape.placeholder_format.idx, shape.placeholder_format.type
    slide = shape.part
    try:
        layout = slide.slide_layout
    except AttributeError:
        return None
    for holder in (layout, layout.slide_master):
        for ph in holder.placeholders:
            if ph.placeholder_format.idx == idx or ph.placeholder_format.type == ph_type:
                sz = from_lststyle(ph.text_frame._txBody, level)
                if sz:
                    return sz
    styles = layout.slide_master._element.find(qn("p:txStyles"))
    if styles is not None:
        name = ("p:titleStyle" if str(ph_type).startswith(("TITLE", "CENTER_TITLE"))
                else "p:bodyStyle")
        st = styles.find(qn(name))
        if st is not None:
            el = st.find(qn(f"a:lvl{level + 1}pPr"))
            if el is not None:
                d = el.find(qn("a:defRPr"))
                if d is not None and d.get("sz"):
                    return int(d.get("sz")) / 100
    return None


def inherited_align(shape, level):
    """A title centred by its layout must render centred, or QA sees collisions
    with the template's own corner art that PowerPoint will never show."""
    if not getattr(shape, "has_text_frame", False):
        return None
    if not shape.is_placeholder:
        return "ctr" if shape.shape_type == 1 else None   # autoshape default
    idx = shape.placeholder_format.idx
    ph_type = shape.placeholder_format.type
    try:
        layout = shape.part.slide_layout
    except AttributeError:
        return None
    for holder in (layout, layout.slide_master):
        for ph in holder.placeholders:
            if ph.placeholder_format.idx == idx or ph.placeholder_format.type == ph_type:
                for para in ph.text_frame.paragraphs:
                    el = para._p.find(qn("a:pPr"))
                    if el is not None and el.get("algn"):
                        return el.get("algn")
                ls = ph.text_frame._txBody.find(qn("a:lstStyle"))
                if ls is not None:
                    lvl = ls.find(qn(f"a:lvl{level + 1}pPr"))
                    if lvl is not None and lvl.get("algn"):
                        return lvl.get("algn")
    styles = layout.slide_master._element.find(qn("p:txStyles"))
    if styles is not None:
        name = "p:titleStyle" if "TITLE" in str(ph_type) else "p:bodyStyle"
        st = styles.find(qn(name))
        if st is not None:
            lvl = st.find(qn(f"a:lvl{level + 1}pPr"))
            if lvl is not None and lvl.get("algn"):
                return lvl.get("algn")
    return None


def wrap(draw, text, fnt, max_px):
    lines, cur = [], ""
    for word in text.split(" "):
        trial = (cur + " " + word).strip()
        if cur and draw.textlength(trial, font=fnt) > max_px:
            lines.append(cur)
            cur = word
        else:
            cur = trial
    lines.append(cur)
    return lines



def para_runs(p, sh, scale, fill, theme):
    """(text, font, colour) for each run of a paragraph, plus explicit breaks."""
    lvl = p.level or 0
    base = inherited_size(sh, lvl) or 18.0
    out = []
    for child in p._p:
        tag = child.tag.split("}")[1]
        if tag == "br":
            out.append(None)
            continue
        if tag != "r":
            continue
        t = child.find(qn("a:t"))
        text = t.text if t is not None and t.text else ""
        rPr = child.find(qn("a:rPr"))
        sz = base
        bold = italic = False
        name = "Arial"
        col = (0, 0, 0)
        if rPr is not None:
            if rPr.get("sz"):
                sz = int(rPr.get("sz")) / 100
            bold = rPr.get("b") == "1"
            italic = rPr.get("i") == "1"
            latin = rPr.find(qn("a:latin"))
            if latin is not None and latin.get("typeface"):
                name = latin.get("typeface")
            solid = rPr.find(qn("a:solidFill"))
            if solid is not None:
                srgb = solid.find(qn("a:srgbClr"))
                sch = solid.find(qn("a:schemeClr"))
                if srgb is not None:
                    v = srgb.get("val")
                    col = tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))
                elif sch is not None:
                    col = theme.rgb(sch.get("val"))
        if fill and sum(fill) < 330 and sum(col) < 200:
            col = (255, 255, 255)
        out.append((text, font(name, sz * scale / 72, bold, italic), col))
    return out


def layout_lines(d, sh, tf, avail, scale, fill):
    """Flow every paragraph's runs into lines the way a renderer would."""
    theme = layout_lines.theme
    lines = []
    for p in tf.paragraphs:
        pPr = p._p.find(qn("a:pPr"))
        algn = pPr.get("algn") if pPr is not None else None
        if algn is None:
            algn = inherited_align(sh, p.level or 0)
        marL = int(pPr.get("marL")) if pPr is not None and pPr.get("marL") else 0
        indent = int(Emu(marL).inches * scale)
        space_after = p.space_after.pt if p.space_after else 0
        bulleted = pPr is not None and pPr.find(qn("a:buChar")) is not None
        runs = para_runs(p, sh, scale, fill, theme)
        width_budget = avail - indent
        cur, curw, curh = [], 0, 0
        def flush(last=False):
            nonlocal cur, curw, curh
            if cur or last:
                lines.append({"runs": list(cur), "width": curw,
                              "h": max(curh, int(12 * scale / 72)) + (
                                  int(space_after * scale / 72) if last else 0),
                              "indent": indent, "align": algn})
            cur, curw, curh = [], 0, 0
        if not runs:
            flush(True)
            continue
        prefix = "• " if bulleted else ""
        for item in runs:
            if item is None:
                flush()
                continue
            text, fnt, col = item
            if prefix:
                text = prefix + text
                prefix = ""
            for word in text.replace("\v", " ").replace("\n", " \x00 ").split(" "):
                if word == "\x00":
                    flush()
                    continue
                piece = word + " "
                pw = d.textlength(piece, font=fnt)
                if cur and curw + pw > width_budget:
                    flush()
                cur.append((piece, fnt, col, curw))
                curw += pw
                curh = max(curh, int(fnt.size * 1.32))
        flush(True)
    return lines


def render(path, outdir, scale=120):
    prs = Presentation(path)
    theme = Theme(prs)
    layout_lines.theme = theme
    W = int(Emu(prs.slide_width).inches * scale)
    H = int(Emu(prs.slide_height).inches * scale)
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)
    problems = []
    made = []

    for n, slide in enumerate(prs.slides, 1):
        bg = solid_fill(slide._element.find(qn("p:cSld")).find(qn("p:bg")), theme) \
            if slide._element.find(qn("p:cSld")).find(qn("p:bg")) is not None else None
        if bg is None:
            bgpr = slide._element.find(f".//{qn('p:bgPr')}")
            bg = solid_fill(bgpr, theme) if bgpr is not None else None
        img = Image.new("RGB", (W, H), bg or (255, 255, 255))
        d = ImageDraw.Draw(img)
        slide_dark = bool(bg) and sum(bg) < 330
        for sh in slide.shapes:
            try:
                x = int(Emu(sh.left or 0).inches * scale)
                y = int(Emu(sh.top or 0).inches * scale)
                w = int(Emu(sh.width or 0).inches * scale)
                h = int(Emu(sh.height or 0).inches * scale)
            except TypeError:
                continue

            if sh.shape_type == 13 and (x + w > W + 2 or y + h > H + 2
                                        or x < -2 or y < -2):
                problems.append(
                    f"slide {n}: picture '{sh.name}' runs off the slide "
                    f"(box {x / scale:.2f},{y / scale:.2f} to "
                    f"{(x + w) / scale:.2f},{(y + h) / scale:.2f}in)")

            if sh.shape_type == 13:  # picture
                try:
                    pic = Image.open(BytesIO(sh.image.blob))
                    pic = pic.convert("RGBA").resize((max(w, 1), max(h, 1)))
                    # Template art is transparent PNG; pasting it opaque turns a
                    # logo into a black rectangle over whatever it sits on, which
                    # then reads as a collision that PowerPoint will never show.
                    img.paste(pic, (x, y), pic)
                except Exception:
                    d.rectangle([x, y, x + w, y + h], fill=(215, 215, 215),
                                outline=(150, 150, 150))
                    d.text((x + 6, y + 6), "[image]", fill=(90, 90, 90),
                           font=font("sans", 14, False, False))
                continue

            if getattr(sh, "has_table", False) and sh.has_table:
                tbl = sh.table
                widths = [int(Emu(c.width).inches * scale) for c in tbl.columns]
                heights = [int(Emu(r.height).inches * scale) for r in tbl.rows]
                cy0 = y
                for ri, rh in enumerate(heights):
                    cx = x
                    for ci, cw in enumerate(widths):
                        cell = tbl.cell(ri, ci)
                        cfill = solid_fill(cell._tc.find(qn("a:tcPr")), theme)
                        d.rectangle([cx, cy0, cx + cw, cy0 + rh],
                                    fill=cfill or (255, 255, 255),
                                    outline=(226, 229, 234))
                        pad2 = int(0.06 * scale)
                        lines = layout_lines(d, sh, cell.text_frame,
                                             cw - 2 * pad2, scale, cfill)
                        ty = cy0 + pad2
                        for ln in lines:
                            for text, fnt, col, dx in ln["runs"]:
                                d.text((cx + pad2 + dx,
                                        ty + ln["h"] - int(fnt.size * 1.15)),
                                       text, fill=col, font=fnt)
                            ty += ln["h"]
                        need = sum(l["h"] for l in lines)
                        if need > rh + 2:
                            problems.append(
                                f"slide {n}: table cell r{ri}c{ci} needs "
                                f"{need / scale:.2f}in in a {rh / scale:.2f}in row")
                        cx += cw
                    cy0 += rh
                continue

            if sh.shape_type == 3:      # chart
                d.rectangle([x, y, x + w, y + h], outline=(200, 205, 212))
                d.text((x + 8, y + 8), "[chart]", fill=(120, 125, 132),
                       font=font("sans", 14, False, False))
                continue

            spPr = sh._element.find(qn("p:spPr"))
            fill = solid_fill(spPr, theme) if spPr is not None else None
            if fill:
                if sh.shape_type == 1 and sh.name.lower().startswith("oval"):
                    d.ellipse([x, y, x + w, y + h], fill=fill)
                else:
                    d.rectangle([x, y, x + w, y + h], fill=fill)

            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            if fill is None and slide_dark:
                fill = bg

            tf = sh.text_frame
            bodyPr = tf._txBody.find(qn("a:bodyPr"))
            anchor = bodyPr.get("anchor") if bodyPr is not None else None
            pad = int(0.05 * scale)
            avail = max(w - 2 * pad, 20)
            lines = layout_lines(d, sh, tf, avail, scale, fill)
            total = sum(ln["h"] for ln in lines)
            if anchor == "ctr":
                cy = y + (h - total) // 2
            elif anchor == "b":
                cy = y + h - total
            else:
                cy = y + pad
            for ln in lines:
                tx = x + pad + ln["indent"]
                if ln["align"] == "ctr":
                    tx = x + (w - ln["width"]) // 2
                elif ln["align"] == "r":
                    tx = x + w - pad - ln["width"]
                for text, fnt, col, dx in ln["runs"]:
                    d.text((tx + dx, cy + ln["h"] - int(fnt.size * 1.15)), text,
                           fill=col, font=fnt)
                cy += ln["h"]

            if total > h + 2:
                problems.append(
                    f"slide {n}: '{sh.name}' (id {sh.shape_id}) text needs "
                    f"{total / scale:.2f}in in a {h / scale:.2f}in box "
                    f"(+{(total - h) / scale:.2f}in)")
            if cy > H:
                problems.append(
                    f"slide {n}: '{sh.name}' (id {sh.shape_id}) text runs "
                    f"{(cy - H) / scale:.2f}in past the bottom of the slide")

        out = outdir / f"slide-{n}.png"
        img.save(out)
        made.append(out)

    print("\n".join(str(p.resolve()) for p in made))
    if problems:
        print("\nfit problems:")
        for p in problems:
            print("  " + p)
    else:
        print("\nno text-fit problems detected")
    return problems


if __name__ == "__main__":
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    scale = 120
    for a in sys.argv[1:]:
        if a.startswith("--scale"):
            scale = int(a.split("=")[1])
    render(args[0], args[1] if len(args) > 1 else "qa", scale)
